"""
FLOP Trade Presentation Generator
===================================
Generates: flop_trade_presentation.pptx (22 slides)
Requires:  pip install python-pptx

Equation images must be in ./eq_imgs/ directory.
Data center photo must be at ./dc_photo.jpg.

Usage:
    python Programs/create_presentation.py
"""

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ============================================================
# COLOR PALETTE
# ============================================================


class C:
    darkBg = RGBColor(0x0B, 0x1D, 0x2E)
    medBg = RGBColor(0x13, 0x2F, 0x46)
    accent = RGBColor(0x00, 0xB4, 0xD8)
    accentDim = RGBColor(0x08, 0x91, 0xB2)
    warm = RGBColor(0xF5, 0x9E, 0x0B)
    lightBg = RGBColor(0xF0, 0xF5, 0xFA)
    cardBg = RGBColor(0xFF, 0xFF, 0xFF)
    textWhite = RGBColor(0xFF, 0xFF, 0xFF)
    textLight = RGBColor(0xCB, 0xD5, 0xE1)
    textDark = RGBColor(0x1E, 0x29, 0x3B)
    textMuted = RGBColor(0x64, 0x74, 0x8B)
    greenAcc = RGBColor(0x10, 0xB9, 0x81)
    redAcc = RGBColor(0xEF, 0x44, 0x44)


HEADER = "Georgia"
BODY = "Calibri"


# ============================================================
# HELPERS
# ============================================================


def hex_to_rgb(h):
    """Convert hex string like '0B1D2E' to RGBColor."""
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, transparency=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if transparency is not None:
        # Set transparency via XML on the shape's spPr
        sp_pr = shape._element.find(qn('p:spPr'))
        if sp_pr is None:
            sp_pr = shape._element
        solid_fill = sp_pr.find(qn('a:solidFill'))
        if solid_fill is not None:
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.makeelement(qn('a:alpha'), {})
                alpha.set('val', str(int((100 - transparency) * 1000)))
                srgb.append(alpha)
    # Remove shadow
    shape.shadow.inherit = False
    return shape


def add_shadow(shape):
    """Apply a subtle outer shadow to a shape."""
    sp = shape._element
    sp_pr = sp.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = sp.find(qn('a:spPr'))
    if sp_pr is None:
        return
    effect_lst = sp_pr.makeelement(qn('a:effectLst'), {})
    outer_shdw = effect_lst.makeelement(qn('a:outerShdw'), {
        'blurRad': '101600',   # 8pt blur
        'dist': '38100',       # 3pt offset
        'dir': '8100000',      # 135 degrees
    })
    srgb_clr = outer_shdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgb_clr.makeelement(qn('a:alpha'), {'val': '12000'})  # 12% opacity
    srgb_clr.append(alpha)
    outer_shdw.append(srgb_clr)
    effect_lst.append(outer_shdw)
    sp_pr.append(effect_lst)


def add_card(slide, x, y, w, h, fill_color=None):
    """Add a card rectangle with shadow."""
    if fill_color is None:
        fill_color = C.cardBg
    shape = add_rect(slide, x, y, w, h, fill_color)
    add_shadow(shape)
    return shape


def add_textbox(slide, x, y, w, h, text, font_name=BODY, font_size=14,
                color=None, bold=False, italic=False, align=PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.TOP, margin_left=0, margin_top=0,
                margin_right=0, margin_bottom=0, word_wrap=True):
    """Add a simple text box with a single run."""
    if color is None:
        color = C.textDark
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.text_frame.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.margin_left = Inches(margin_left) if margin_left else Emu(0)
    tf.margin_top = Inches(margin_top) if margin_top else Emu(0)
    tf.margin_right = Inches(margin_right) if margin_right else Emu(0)
    tf.margin_bottom = Inches(margin_bottom) if margin_bottom else Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Set vertical alignment
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        bodyPr.set('anchor', anchor_map.get(valign, 't'))

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_rich_text(slide, x, y, w, h, runs_data, valign=MSO_ANCHOR.TOP):
    """
    Add a text box with multiple styled runs.
    runs_data: list of dicts with keys: text, font_name, font_size, color, bold, italic, break_line, space_after
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_bottom = Emu(0)

    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
        bodyPr.set('anchor', anchor_map.get(valign, 't'))

    current_para = tf.paragraphs[0]
    current_para.space_before = Pt(0)
    current_para.space_after = Pt(0)

    for rd in runs_data:
        if rd.get('new_para', False):
            current_para = tf.add_paragraph()
            current_para.space_before = Pt(0)
            sa = rd.get('space_after', 0)
            current_para.space_after = Pt(sa) if sa else Pt(0)
            if rd.get('align'):
                current_para.alignment = rd['align']

        run = current_para.add_run()
        run.text = rd.get('text', '')
        run.font.name = rd.get('font_name', BODY)
        run.font.size = Pt(rd.get('font_size', 12))
        run.font.color.rgb = rd.get('color', C.textDark)
        run.font.bold = rd.get('bold', False)
        run.font.italic = rd.get('italic', False)

        if rd.get('break_line', False):
            current_para = tf.add_paragraph()
            sa = rd.get('space_after', 0)
            current_para.space_after = Pt(sa) if sa else Pt(0)
            current_para.space_before = Pt(0)

    return txBox


def add_oval(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_image_safe(slide, path, x, y, w, h):
    """Add image if file exists, otherwise add a placeholder rect."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        add_rect(slide, x, y, w, h, hex_to_rgb("334455"))
        add_textbox(slide, x, y, w, h, f"[Image: {path}]",
                    font_size=9, color=C.textMuted, align=PP_ALIGN.CENTER,
                    valign=MSO_ANCHOR.MIDDLE)


def regime_color(rtype):
    if "EE" in rtype:
        return C.greenAcc
    if "IE" in rtype:
        return C.accent
    if "DD" in rtype:
        return C.warm
    return C.textMuted


# ============================================================
# MAIN
# ============================================================


def main():
    # Resolve asset paths relative to project root (one level up from Programs/)
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    os.chdir(project_root)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ============================================================
    # SLIDE 1: TITLE
    # ============================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, C.darkBg)

    # Top accent bar
    add_rect(s1, 0, 0, 10, 0.06, C.accent)

    # Title
    add_textbox(s1, 0.8, 1.2, 8, 2.0, "Cheap Energy\nMight Not Be Enough",
                font_name=HEADER, font_size=40, color=C.textWhite, bold=True)

    # Subtitle
    add_textbox(s1, 0.8, 3.1, 8, 0.6, "A Trade Model of AI Compute Services",
                font_name=BODY, font_size=20, color=C.accent)

    # Divider
    add_rect(s1, 0.8, 3.9, 2.5, 0.04, C.accent)

    # Author
    add_textbox(s1, 0.8, 4.15, 5, 0.45, "Michael Lokshin",
                font_name=BODY, font_size=16, color=C.textLight)
    add_textbox(s1, 0.8, 4.55, 5, 0.4, "March 2026",
                font_name=BODY, font_size=14, color=C.textMuted)

    # ============================================================
    # SLIDE 2: DATA CENTER — HOW FLOPS ARE PRODUCED
    # ============================================================
    s1x = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1x, C.darkBg)

    # Data center photo
    add_image_safe(s1x, "dc_photo.jpg", 0, 0, 10, 3.15)

    # Dark overlay
    add_rect(s1x, 0, 0, 10, 1.2, RGBColor(0, 0, 0), transparency=35)
    add_rect(s1x, 0, 2.5, 10, 0.65, C.darkBg, transparency=25)

    # Title over photo
    add_textbox(s1x, 0.5, 0.15, 9.2, 0.55, "How FLOPs Are Produced",
                font_name=HEADER, font_size=30, color=C.textWhite, bold=True)
    add_textbox(s1x, 0.5, 0.65, 7.5, 0.35,
                "A data center converts electricity into floating-point operations — the raw commodity of AI",
                font_name=BODY, font_size=12, color=hex_to_rgb("CFD8DC"), italic=True)

    # Production pipeline: 4 stages below photo
    stages = [
        {"label": "Electricity", "sub": "Grid power (pⱼᴱ)\nLocal price — only\ncountry-varying input", "color": "00B4D8", "icon": "⚡"},
        {"label": "Cooling", "sub": "PUE overhead\nClimate-dependent\n(θⱼ → cooling cost)", "color": "0097A7", "icon": "🏭"},
        {"label": "GPU Hardware", "sub": "H100/B200 racks\nGlobally priced\n(~90% of unit cost)", "color": "00ACC1", "icon": "🖥"},
        {"label": "FLOPs Output", "sub": "Compute services\nTraining (batch)\nInference (real-time)", "color": "F59E0B", "icon": "📊"},
    ]

    for i, st in enumerate(stages):
        xPos = 0.3 + i * 2.45
        yPos = 3.25
        clr = hex_to_rgb(st["color"])

        # Card
        add_card(s1x, xPos, yPos, 2.2, 1.35, hex_to_rgb("1A2332"))
        # Top accent
        add_rect(s1x, xPos, yPos, 2.2, 0.05, clr)
        # Icon
        add_textbox(s1x, xPos + 0.1, yPos + 0.15, 0.35, 0.35, st["icon"],
                    font_size=18, color=clr, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Label
        add_textbox(s1x, xPos + 0.5, yPos + 0.12, 1.55, 0.3, st["label"],
                    font_size=13, color=C.textWhite, bold=True)
        # Sub text
        add_textbox(s1x, xPos + 0.15, yPos + 0.5, 1.9, 0.75, st["sub"],
                    font_size=9, color=C.textMuted)
        # Arrow
        if i < 3:
            add_textbox(s1x, xPos + 2.15, yPos + 0.2, 0.35, 0.35, "→",
                        font_name=HEADER, font_size=20, color=C.accent,
                        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Bottom strip: cost equation
    add_rect(s1x, 0, 4.7, 10, 0.55, hex_to_rgb("111B27"))
    add_rich_text(s1x, 0.5, 4.72, 9.0, 0.5, [
        {"text": "Unit cost:  ", "color": C.warm, "bold": True, "font_size": 12},
        {"text": "cⱼ = ", "color": C.textWhite, "font_size": 12},
        {"text": "electricity (~8%)", "color": C.accent, "bold": True, "font_size": 12},
        {"text": "  +  ", "color": C.textWhite, "font_size": 12},
        {"text": "hardware (~90%)", "color": hex_to_rgb("00ACC1"), "bold": True, "font_size": 12},
        {"text": "  +  ", "color": C.textWhite, "font_size": 12},
        {"text": "networking + construction (~2%)", "color": C.textMuted, "bold": True, "font_size": 12},
    ], valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 3: GLOBAL MEGAPROJECTS — DEVELOPING COUNTRIES
    # ============================================================
    s1y = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1y, C.darkBg)

    add_textbox(s1y, 0.5, 0.1, 9.2, 0.5, "Developing Countries Are Joining the Race",
                font_name=HEADER, font_size=26, color=C.textWhite, bold=True)
    add_textbox(s1y, 0.5, 0.55, 9.2, 0.3,
                "Billion-dollar data center commitments are no longer confined to the US, Europe, and China",
                font_size=11, color=C.textMuted, italic=True)

    megaprojects = [
        {"region": "ARMENIA", "project": "Firebird / NVIDIA AI Factory",
         "detail": "$4B total (Phase 1: $500M) · 100 MW\n50,000 NVIDIA GPUs · US-approved chip transfer\nPM Pashinyan lobbied Trump directly",
         "color": "F59E0B"},
        {"region": "SAUDI ARABIA", "project": "HUMAIN / NEOM / DataVolt",
         "detail": "$23B+ in AI partnerships · 11 data centers\n200 MW each · NEOM pivoting to DC hub\nPIF goal: 3rd-largest AI provider globally",
         "color": "00BCD4"},
        {"region": "UAE", "project": "OpenAI–G42 Abu Dhabi",
         "detail": "5 GW campus · largest single AI site globally\nStargate international expansion\nKhazna 100 MW AI facility in Ajman",
         "color": "00B4D8"},
        {"region": "INDONESIA", "project": "Jakarta / Batam DC corridor",
         "detail": "Princeton Digital: $1B Jakarta campus\nMicrosoft & Google cloud regions\nLargest SE Asian market by demand",
         "color": "4CAF50"},
        {"region": "KENYA", "project": "Africa Data Centres / IXAfrica",
         "detail": "East Africa's DC hub · 1 GW pipeline\nCheap geothermal power (~$0.07/kWh)\nMicrosoft & Google African cloud regions",
         "color": "E91E63"},
        {"region": "MALAYSIA", "project": "Johor corridor (Microsoft / NTT)",
         "detail": "$34B investment boom · $6B market by 2031\nMicrosoft 2nd cloud region (SE Asia 3)\nCyberjaya: 22 existing + 9 upcoming DCs",
         "color": "7C4DFF"},
    ]

    for i, mp in enumerate(megaprojects):
        col = i % 3
        row = i // 3
        xPos = 0.35 + col * 3.15
        yPos = 1.0 + row * 1.75
        clr = hex_to_rgb(mp["color"])

        # Card
        add_card(s1y, xPos, yPos, 3.0, 1.55, hex_to_rgb("1A2332"))
        # Left accent bar
        add_rect(s1y, xPos, yPos, 0.06, 1.55, clr)
        # Region label
        add_textbox(s1y, xPos + 0.2, yPos + 0.08, 2.6, 0.25, mp["region"],
                    font_size=9, color=clr, bold=True)
        # Project name
        add_textbox(s1y, xPos + 0.2, yPos + 0.32, 2.6, 0.32, mp["project"],
                    font_size=11, color=C.textWhite, bold=True)
        # Detail
        add_textbox(s1y, xPos + 0.2, yPos + 0.68, 2.6, 0.8, mp["detail"],
                    font_size=9, color=C.textMuted)

    # Bottom bar
    add_rect(s1y, 0, 4.65, 10, 0.6, hex_to_rgb("111B27"))
    add_rich_text(s1y, 0.5, 4.65, 9.0, 0.6, [
        {"text": "Key question: ", "color": C.warm, "bold": True, "font_size": 12},
        {"text": "These countries have cheap energy and political will. Can they overcome sovereignty frictions and institutional gaps to capture AI compute demand?",
         "color": C.textLight, "font_size": 12},
    ], valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 4: MOTIVATION — THE OPPORTUNITY
    # ============================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, C.lightBg)

    add_textbox(s2, 0.7, 0.4, 9, 0.7, "The Opportunity: FLOP Exporting",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)

    # Left column — narrative card
    add_card(s2, 0.7, 1.4, 5.3, 3.6)
    add_rect(s2, 0.7, 1.4, 0.07, 3.6, C.accent)

    paras = [
        "AI compute demand is doubling every six months. Data centers already consume 1.5% of global electricity — projected to more than double by 2030.",
        "Countries can convert cheap electricity into high-value AI compute exports — moving up the value chain from raw resources to digital services.",
        "Energy-rich developing countries — with cheap electricity, growing grid capacity, and political will — have a rare opportunity to leapfrog into high-value digital exports without the traditional industrialization path.",
    ]
    y_off = 1.6
    for para in paras:
        add_textbox(s2, 1.1, y_off, 4.6, 0.9, para, font_size=14, color=C.textDark)
        y_off += 1.0

    # Right column — 3 stat callouts
    stats = [
        {"num": "2×", "label": "compute demand\ndoubling rate", "icon": "⚡"},
        {"num": "85", "label": "countries\ncalibrated", "icon": "🌍"},
        {"num": "$9B+", "label": "cloud exports\nannually", "icon": "🖥"},
    ]
    for i, st in enumerate(stats):
        yPos = 1.4 + i * 1.25
        add_card(s2, 6.4, yPos, 3.0, 1.05)
        add_textbox(s2, 6.65, yPos + 0.22, 0.45, 0.45, st["icon"],
                    font_size=20, color=C.accent, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(s2, 7.25, yPos + 0.05, 1.5, 0.55, st["num"],
                    font_name=HEADER, font_size=28, color=C.accent, bold=True)
        add_textbox(s2, 7.25, yPos + 0.55, 2.0, 0.45, st["label"],
                    font_size=11, color=C.textMuted)

    # ============================================================
    # SLIDE 5: PAPER CONTRIBUTIONS
    # ============================================================
    s1c = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1c, C.lightBg)

    add_textbox(s1c, 0.5, 0.25, 9.2, 0.5, "This Paper's Contributions",
                font_name=HEADER, font_size=26, color=C.textDark, bold=True)

    contribs = [
        {"num": "1", "title": "Model",
         "desc": "Capacity-constrained Ricardian model of compute trade with iceberg latency costs and bilateral sovereignty premia",
         "detail": "Distinguishes training (latency-insensitive) from inference (latency-sensitive) — five equilibrium regimes emerge"},
        {"num": "2", "title": "Calibration",
         "desc": "85-country calibration correcting for energy subsidies that distort headline cost rankings",
         "detail": "Cost-recovery pricing replaces subsidized tariffs with LRMC for 13 countries; bilateral sovereignty premium from UNGA voting, data adequacy, and sanctions"},
        {"num": "3", "title": "Welfare & Policy",
         "desc": "Characterize trade regimes and quantify the welfare cost of sovereignty premia",
         "detail": "Welfare cost is 1.5% of compute spending (~$1.3B/yr) — small enough that domestic production can be rational for sovereign workloads"},
    ]

    for i, c in enumerate(contribs):
        yPos = 1.05 + i * 1.3
        # Card
        add_card(s1c, 1.2, yPos, 7.8, 1.1)
        add_rect(s1c, 1.2, yPos, 0.07, 1.1, C.accent)
        # Number circle
        add_oval(s1c, 0.35, yPos + 0.2, 0.65, 0.65, C.accent)
        add_textbox(s1c, 0.35, yPos + 0.2, 0.65, 0.65, c["num"],
                    font_name=HEADER, font_size=24, color=C.textWhite, bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Title
        add_textbox(s1c, 2.05, yPos + 0.1, 2.0, 0.4, c["title"],
                    font_size=15, color=C.textDark, bold=True)
        # Desc
        add_textbox(s1c, 2.05, yPos + 0.45, 6.7, 0.3, c["desc"],
                    font_size=12, color=C.textDark)
        # Detail
        add_textbox(s1c, 2.05, yPos + 0.72, 6.7, 0.3, c["detail"],
                    font_size=10, color=C.textMuted)

    # ============================================================
    # SLIDE 6: RELATED LITERATURE
    # ============================================================
    s1b = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1b, C.lightBg)

    add_textbox(s1b, 0.5, 0.25, 9.2, 0.5, "Related Literature",
                font_name=HEADER, font_size=26, color=C.textDark, bold=True)
    add_textbox(s1b, 0.5, 0.7, 9.2, 0.35,
                "Existing work examines AI governance, digital trade, and data center geography — but no formal trade model of compute exists",
                font_size=12, color=C.textMuted, italic=True)

    litCards = [
        {"title": "AI & Comparative Advantage",
         "refs": "Goldfarb & Trefler (2018)\nKorinek & Stiglitz (2021)",
         "note": "AI shifts advantage to data & human capital; developing countries risk being left behind",
         "x": 0.5, "y": 1.2},
        {"title": "IT Services Offshoring",
         "refs": "Blinder (2006)\nBrainard (1997), Helpman et al. (2004)",
         "note": "Labor-intensive and skill-biased — a different set of exporters than compute",
         "x": 5.1, "y": 1.2},
        {"title": "Data Center Geography",
         "refs": "Flucker et al. (2013), Liu et al. (2023)\nLehdonvirta et al. (2024), Pilz et al. (2025)",
         "note": "Climate affects cooling costs; 'Compute North' holds 77% of capacity",
         "x": 0.5, "y": 3.0},
        {"title": "Compute Governance",
         "refs": "Sastry et al. (2024)\nWorld Bank (2025), Biglaiser et al. (2024)",
         "note": "Compute is well-suited for regulation but lacks a cost-to-trade framework",
         "x": 5.1, "y": 3.0},
    ]

    for lc in litCards:
        add_card(s1b, lc["x"], lc["y"], 4.4, 1.6)
        add_rect(s1b, lc["x"], lc["y"], 4.4, 0.06, C.accent)
        add_textbox(s1b, lc["x"] + 0.2, lc["y"] + 0.15, 4.0, 0.35, lc["title"],
                    font_size=13, color=C.textDark, bold=True)
        add_textbox(s1b, lc["x"] + 0.2, lc["y"] + 0.48, 4.0, 0.45, lc["refs"],
                    font_size=10, color=C.accent)
        add_textbox(s1b, lc["x"] + 0.2, lc["y"] + 1.0, 4.0, 0.5, lc["note"],
                    font_size=11, color=C.textMuted)

    # Gap callout
    add_rect(s1b, 0.5, 4.75, 9.0, 0.5, C.darkBg)
    add_textbox(s1b, 1.15, 4.75, 8.2, 0.5,
                "Gap: No formal model linking production costs → trade patterns → welfare costs for compute services",
                font_size=12, color=C.textLight, bold=True, valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 7: MODEL OVERVIEW
    # ============================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, C.lightBg)

    add_textbox(s3, 0.7, 0.4, 9, 0.7, "Model Architecture",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)
    add_textbox(s3, 0.7, 1.0, 9, 0.4, "A capacity-constrained Ricardian model with three innovations",
                font_size=15, color=C.textMuted, italic=True)

    pillars = [
        {"title": "Cost Structure", "desc": "Country-specific production costs from electricity, hardware (globally priced), and construction. PUE varies with climate.", "color": C.accent},
        {"title": "Trade Costs", "desc": "Training is latency-insensitive (freely offshorable). Inference degrades with distance — localized markets.", "color": C.accentDim},
        {"title": "Sovereignty Premium", "desc": "Bilateral friction from geopolitical distance, regulatory gaps, and sanctions. Trust deficits as trade barriers.", "color": C.warm},
    ]

    for i, p in enumerate(pillars):
        xPos = 0.7 + i * 3.1
        add_card(s3, xPos, 1.7, 2.8, 3.3)
        add_rect(s3, xPos, 1.7, 2.8, 0.06, p["color"])
        add_textbox(s3, xPos + 0.3, 2.6, 2.2, 0.45, p["title"],
                    font_name=HEADER, font_size=16, color=C.textDark, bold=True)
        add_textbox(s3, xPos + 0.3, 3.1, 2.2, 1.6, p["desc"],
                    font_size=12, color=C.textMuted)

    # ============================================================
    # SLIDE 8: FORMAL MODEL — COST & TRADE EQUATIONS
    # ============================================================
    s3b = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3b, C.lightBg)

    add_textbox(s3b, 0.5, 0.2, 9.2, 0.45, "Model: Cost Structure & Trade Costs",
                font_name=HEADER, font_size=25, color=C.textDark, bold=True)

    # Equation 1: Production cost
    add_card(s3b, 0.5, 0.8, 9.0, 1.55)
    add_rect(s3b, 0.5, 0.8, 0.07, 1.55, C.accent)
    add_textbox(s3b, 0.8, 0.85, 3, 0.3, "(1)  Production Cost",
                font_size=13, color=C.accent, bold=True)
    add_image_safe(s3b, "eq_imgs/eq1.png", 0.9, 1.1, 4.5, 0.5)
    add_rich_text(s3b, 5.5, 0.82, 3.8, 1.5, [
        {"text": "γ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = GPU power draw (kW)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "pⱼᴱ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = electricity price ($/kWh)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "PUE(θⱼ)", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = cooling overhead (θ = temp.)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "ρ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = hardware amortization (global, ~90%)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "η", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = networking (global)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "κⱼ/(Lf·H)", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = construction / (lifetime × hrs/yr)", "color": C.textMuted, "font_size": 10},
    ])

    # Equation 2: Sovereignty premium
    add_card(s3b, 0.5, 2.5, 9.0, 1.35)
    add_rect(s3b, 0.5, 2.5, 0.07, 1.35, C.warm)
    add_textbox(s3b, 0.8, 2.55, 4, 0.3, "(2)  Bilateral Sovereignty Premium",
                font_size=13, color=C.warm, bold=True)
    add_image_safe(s3b, "eq_imgs/eq2.png", 0.9, 2.82, 4.5, 0.5)
    add_rich_text(s3b, 5.5, 2.55, 3.8, 1.2, [
        {"text": "dⱼₖ", "color": C.warm, "bold": True, "font_size": 10},
        {"text": " = geopolitical distance (UNGA voting)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "Rⱼₖ", "color": C.warm, "bold": True, "font_size": 10},
        {"text": " = data-adequacy agreement (0/1)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "Sⱼₖ", "color": C.warm, "bold": True, "font_size": 10},
        {"text": " = sanctions indicator → λ = ∞", "color": C.textMuted, "font_size": 10},
    ])

    # Equation 3: Delivered cost
    add_card(s3b, 0.5, 4.0, 9.0, 1.3)
    add_rect(s3b, 0.5, 4.0, 0.07, 1.3, C.accent)
    add_textbox(s3b, 0.8, 4.05, 3, 0.3, "(3)  Delivered Cost",
                font_size=13, color=C.accent, bold=True)
    add_image_safe(s3b, "eq_imgs/eq3.png", 0.9, 4.32, 5.0, 0.5)
    add_rich_text(s3b, 6.0, 3.98, 3.3, 1.25, [
        {"text": "τ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = latency degradation rate", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "ℓⱼₖ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = round-trip latency (ms)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 4},
        {"text": "Training (s=T):", "color": C.accent, "bold": True, "font_size": 10, "break_line": True},
        {"text": "  τ = 0 → no latency cost", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 2},
        {"text": "Inference (s=I):", "color": C.accent, "bold": True, "font_size": 10, "break_line": True},
        {"text": "  τ > 0 → cost rises with distance", "color": C.textMuted, "font_size": 10},
    ])

    # ============================================================
    # SLIDE 9: EQUILIBRIUM & WELFARE
    # ============================================================
    s3c = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3c, C.lightBg)

    add_textbox(s3c, 0.5, 0.2, 9.2, 0.45, "Model: Equilibrium & Welfare",
                font_name=HEADER, font_size=25, color=C.textDark, bold=True)

    # Equation 4: Demand
    add_card(s3c, 0.5, 0.8, 9.0, 1.1)
    add_rect(s3c, 0.5, 0.8, 0.07, 1.1, C.accent)
    add_textbox(s3c, 0.8, 0.85, 3, 0.3, "(4)  Demand",
                font_size=13, color=C.accent, bold=True)
    add_image_safe(s3c, "eq_imgs/eq4.png", 0.9, 1.08, 4.5, 0.5)
    add_rich_text(s3c, 5.5, 0.85, 3.8, 0.9, [
        {"text": "ωₖ", "color": C.accent, "bold": True, "font_size": 10},
        {"text": " = share of installed DC capacity (MW)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "Splits: Dᵀ = αΩωₖ  (training),  Dᴵ = (1−α)Ωωₖ  (inference)", "color": C.textMuted, "font_size": 10},
    ])

    # Equation 5: Sourcing
    add_card(s3c, 0.5, 2.05, 9.0, 1.35)
    add_rect(s3c, 0.5, 2.05, 0.07, 1.35, C.accent)
    add_textbox(s3c, 0.8, 2.1, 5, 0.3, "(5)  Sourcing & Capacity Constraint",
                font_size=13, color=C.accent, bold=True)
    add_image_safe(s3c, "eq_imgs/eq5.png", 0.9, 2.4, 5.0, 0.5)
    add_rich_text(s3c, 5.5, 2.1, 3.8, 1.2, [
        {"text": "Each buyer picks cheapest delivered source", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "GPU-hours allocated to highest-margin use", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "Subject to grid capacity ceiling k̄ⱼ", "color": C.textMuted, "font_size": 10},
    ])

    # Equation 6: Training market clearing
    add_card(s3c, 0.5, 3.55, 9.0, 1.0)
    add_rect(s3c, 0.5, 3.55, 0.07, 1.0, C.accent)
    add_textbox(s3c, 0.8, 3.58, 5, 0.3, "(6)  Training Market Clearing",
                font_size=13, color=C.accent, bold=True)
    add_image_safe(s3c, "eq_imgs/eq6.png", 0.9, 3.82, 5.0, 0.5)
    add_rich_text(s3c, 6.0, 3.6, 3.3, 0.85, [
        {"text": "Price = marginal exporter's cost", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 3},
        {"text": "Infra-marginal exporters earn rents", "color": C.textMuted, "font_size": 10},
    ])

    # Welfare bar
    add_rect(s3c, 0.5, 4.55, 9.0, 0.75, C.darkBg)
    add_textbox(s3c, 0.7, 4.57, 1.0, 0.3, "Welfare:",
                font_size=11, color=C.warm, bold=True)
    add_image_safe(s3c, "eq_imgs/eq_welfare_white.png", 1.5, 4.57, 4.5, 0.55)
    add_textbox(s3c, 6.1, 4.55, 3.2, 0.75,
                "Import markup + allocative inefficiency\nfrom domestic production",
                font_size=10, color=C.textLight, valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 10: DEMAND & CAPACITY CONSTRAINTS
    # ============================================================
    s3d = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3d, C.lightBg)

    add_textbox(s3d, 0.5, 0.2, 9.2, 0.45, "Demand & Capacity Constraints",
                font_name=HEADER, font_size=25, color=C.textDark, bold=True)

    # LEFT PANEL: Demand
    add_card(s3d, 0.5, 0.8, 4.3, 4.15)
    add_rect(s3d, 0.5, 0.8, 4.3, 0.06, C.accent)
    add_textbox(s3d, 0.75, 0.9, 3.8, 0.35, "How demand is defined",
                font_size=15, color=C.accent, bold=True)

    demand_steps = [
        {"num": "1", "y": 1.4, "title": "Global spending Ω", "sub": "Total cloud compute market (~$90B/yr)"},
        {"num": "2", "y": 2.05, "title": "Country shares  ωₖ = k̄ₖ / Σk̄ₖ'", "sub": "Proxied by installed DC capacity (MW)\nUS ~43%, China ~26%, rest of world ~31%"},
        {"num": "3", "y": 2.85, "title": "Training / inference split", "sub": "Dᵀ = α · Ω · ωₖ     (latency-insensitive)\nDᴵ  = (1−α) · Ω · ωₖ  (latency-sensitive)\nBaseline: α = 50% training, 50% inference"},
    ]
    for ds in demand_steps:
        add_oval(s3d, 0.75, ds["y"], 0.4, 0.4, C.accent)
        add_textbox(s3d, 0.75, ds["y"], 0.4, 0.4, ds["num"],
                    font_name=HEADER, font_size=14, color=C.textWhite, bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(s3d, 1.3, ds["y"] - 0.05, 3.3, 0.25, ds["title"],
                    font_size=12, color=C.textDark, bold=True)
        add_textbox(s3d, 1.3, ds["y"] + 0.22, 3.3, 0.6, ds["sub"],
                    font_size=10, color=C.textMuted)

    add_rich_text(s3d, 0.75, 3.85, 3.85, 0.9, [
        {"text": "Key: ", "color": C.accent, "bold": True, "font_size": 11},
        {"text": "demand shares are fixed — the model determines who supplies each country's demand, not how much each country demands.",
         "color": C.textMuted, "font_size": 11},
    ])

    # RIGHT PANEL: Capacity constraints
    add_card(s3d, 5.1, 0.8, 4.4, 4.15)
    add_rect(s3d, 5.1, 0.8, 4.4, 0.06, C.warm)
    add_textbox(s3d, 5.35, 0.9, 4.0, 0.35, "How capacity constraints work",
                font_size=15, color=C.warm, bold=True)

    add_rich_text(s3d, 5.35, 1.3, 4.0, 1.35, [
        {"text": "Each country has a capacity ceiling k̄ⱼ ", "color": C.textDark, "bold": True, "font_size": 11, "break_line": True},
        {"text": "(grid power + permitting + connectivity)", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 6},
        {"text": "Without constraints: ", "color": C.textDark, "font_size": 10},
        {"text": "cheapest country serves all demand, pᵀ = c₁, zero rents", "color": C.textMuted, "font_size": 10, "break_line": True, "space_after": 4},
        {"text": "With constraints: ", "color": C.textDark, "font_size": 10},
        {"text": "cheapest fills up → price rises to marginal entrant's cost → infra-marginal exporters earn rents",
         "color": C.textMuted, "font_size": 10},
    ])

    # Supply stack bars
    bars = [
        {"label": "KGZ", "w": 0.6, "cost": "$1.58", "color": "4DD0E1", "h_frac": 0.61},
        {"label": "CAN", "w": 1.1, "cost": "$1.59", "color": "26C6DA", "h_frac": 0.71},
        {"label": "KOS", "w": 0.45, "cost": "$1.60", "color": "00BCD4", "h_frac": 0.79},
        {"label": "FIN", "w": 0.7, "cost": "$1.61", "color": "00ACC1", "h_frac": 0.86},
        {"label": "...", "w": 0.55, "cost": "$1.65+", "color": "0097A7", "h_frac": 1.0},
    ]
    barX = 5.4
    barBaseY = 3.0
    maxH = 1.4
    for b in bars:
        h = b["h_frac"] * maxH
        y = barBaseY + (maxH - h)
        add_rect(s3d, barX, y, b["w"], h, hex_to_rgb(b["color"]))
        add_textbox(s3d, barX, barBaseY + maxH - 0.25, b["w"], 0.22, b["label"],
                    font_size=7, color=C.textWhite, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s3d, barX, y - 0.18, b["w"], 0.18, b["cost"],
                    font_size=7, color=C.textDark, bold=True, align=PP_ALIGN.CENTER)
        barX += b["w"] + 0.06

    add_textbox(s3d, 5.4, barBaseY + maxH + 0.05, 3.5, 0.2,
                "← cumulative capacity (GPU-hours) →",
                font_size=8, color=C.textMuted, italic=True, align=PP_ALIGN.CENTER)

    # Supply stack annotations: price line (thin rect) and rents
    add_rect(s3d, 5.3, barBaseY - 0.01, barX - 5.3 + 0.1, 0.02, C.warm)
    add_textbox(s3d, barX + 0.15, barBaseY - 0.15, 1.8, 0.3, "pᵀ = marginal entrant",
                font_size=8, color=C.warm, bold=True)
    add_textbox(s3d, 5.05, barBaseY + 0.3, 0.35, 0.6, "↕\nrents",
                font_size=8, color=C.warm, bold=True, align=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 11: COST STRUCTURE — THE COMPRESSION
    # ============================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, C.lightBg)

    add_textbox(s4, 0.7, 0.4, 9, 0.7, "Why the Cost Spread Is Narrow",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)

    # Stacked bar chart
    chart_data = {
        "categories": ["Cheapest", "Median", "Costliest"],
        "series": [
            {"name": "Hardware", "values": [1.36, 1.36, 1.36]},
            {"name": "Electricity", "values": [0.06, 0.12, 0.22]},
            {"name": "Construction", "values": [0.04, 0.06, 0.10]},
        ]
    }
    cd = CategoryChartData()
    cd.categories = chart_data["categories"]
    for s in chart_data["series"]:
        cd.add_series(s["name"], s["values"])
    chart_frame = s4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.8), cd
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    # Color series
    colors_chart = [C.accent, C.warm, C.textMuted]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors_chart[i]

    # Key insight card
    add_card(s4, 6.4, 1.3, 3.2, 2.0, C.darkBg)
    add_textbox(s4, 6.4, 1.5, 3.2, 0.8, "~90%",
                font_name=HEADER, font_size=48, color=C.accent, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s4, 6.6, 2.3, 2.8, 0.8, "of unit cost is hardware\n— identical across countries",
                font_size=13, color=C.textLight, align=PP_ALIGN.CENTER)

    # Second insight
    add_card(s4, 6.4, 3.6, 3.2, 1.5)
    add_rect(s4, 6.4, 3.6, 0.07, 1.5, C.warm)
    add_textbox(s4, 6.7, 3.7, 2.8, 0.55, "12–20%",
                font_name=HEADER, font_size=30, color=C.warm, bold=True)
    add_textbox(s4, 6.7, 4.25, 2.8, 0.75,
                "total cross-country cost spread\n— narrower than virtually any\nother tradable sector",
                font_size=12, color=C.textMuted)

    # ============================================================
    # SLIDE 12: TRAINING vs INFERENCE
    # ============================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, C.lightBg)

    add_textbox(s5, 0.7, 0.4, 9, 0.7, "Two Markets, Two Geographies",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)

    # Training card
    add_card(s5, 0.7, 1.4, 4.15, 3.6)
    add_rect(s5, 0.7, 1.4, 4.15, 0.06, C.accent)
    add_textbox(s5, 1.0, 1.7, 3.5, 0.4, "TRAINING",
                font_size=14, color=C.accent, bold=True)
    add_textbox(s5, 1.0, 2.1, 3.5, 0.35, "Latency-Insensitive",
                font_name=HEADER, font_size=18, color=C.textDark, bold=True)
    training_points = "• Model training, fine-tuning, batch processing\n• No distance penalty — freely offshorable\n• Global market: cheapest producer wins\n• Weeks to months of compute time\n• Highly concentrated: 1–2 exporters serve most demand"
    add_textbox(s5, 1.0, 2.6, 3.6, 2.2, training_points,
                font_size=13, color=C.textMuted)

    # Inference card
    add_card(s5, 5.15, 1.4, 4.15, 3.6)
    add_rect(s5, 5.15, 1.4, 4.15, 0.06, C.warm)
    add_textbox(s5, 5.45, 1.7, 3.5, 0.4, "INFERENCE",
                font_size=14, color=C.warm, bold=True)
    add_textbox(s5, 5.45, 2.1, 3.5, 0.35, "Latency-Sensitive",
                font_name=HEADER, font_size=18, color=C.textDark, bold=True)
    inference_points = "• Chatbots, autonomous systems, real-time agents\n• Quality degrades with distance (latency)\n• Regional market: proximity matters\n• Millisecond response required\n• More dispersed: regional hubs serve nearby demand"
    add_textbox(s5, 5.45, 2.6, 3.6, 2.2, inference_points,
                font_size=13, color=C.textMuted)

    # ============================================================
    # SLIDE 13: FIVE EQUILIBRIUM REGIMES
    # ============================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, C.lightBg)

    add_textbox(s6, 0.7, 0.4, 9, 0.7, "Five Equilibrium Regimes",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)
    add_textbox(s6, 0.7, 1.0, 9, 0.35,
                "Of nine possible combinations, only five arise in equilibrium (Proposition 1)",
                font_size=14, color=C.textMuted, italic=True)

    regimes = [
        {"code": "EE", "train": "Export", "infer": "Export", "desc": "Cheapest producers — global training + regional inference", "clr": C.greenAcc},
        {"code": "IE", "train": "Import", "infer": "Export", "desc": "Regional inference hubs — not cheap enough for global training", "clr": C.accent},
        {"code": "ID", "train": "Import", "infer": "Domestic", "desc": "Import training, produce inference domestically", "clr": C.accentDim},
        {"code": "DD", "train": "Domestic", "infer": "Domestic", "desc": "High sovereignty premium — produce both domestically", "clr": C.warm},
        {"code": "II", "train": "Import", "infer": "Import", "desc": "High-cost countries — import both services", "clr": C.textMuted},
    ]

    # Build table
    tbl_shape = s6.shapes.add_table(len(regimes) + 1, 4, Inches(0.7), Inches(1.6), Inches(8.6), Inches(3.0))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.0)
    tbl.columns[1].width = Inches(1.2)
    tbl.columns[2].width = Inches(1.2)
    tbl.columns[3].width = Inches(5.2)

    headers = ["Regime", "Training", "Inference", "Description"]
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = C.textWhite
                r.font.name = BODY
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.darkBg

    for i, reg in enumerate(regimes):
        row_bg = hex_to_rgb("F8FAFC") if i % 2 == 0 else C.cardBg
        data = [reg["code"], reg["train"], reg["infer"], reg["desc"]]
        for j, val in enumerate(data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j < 3 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.name = BODY
                    r.font.color.rgb = reg["clr"] if j == 0 else (C.textDark if j < 3 else C.textMuted)
                    r.font.bold = (j == 0)

    add_textbox(s6, 0.7, 4.9, 8.6, 0.35,
                "The remaining four combinations (e.g., training exporter + inference importer) are ruled out by cost ordering.",
                font_size=11, color=C.textMuted, italic=True)

    # ============================================================
    # SLIDE 14: PROPOSITION 4 — NESTING RESULT
    # ============================================================
    s6b = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6b, C.lightBg)

    add_textbox(s6b, 0.5, 0.2, 9.2, 0.45, "Training Exporters ⊂ Inference Exporters",
                font_name=HEADER, font_size=26, color=C.textDark, bold=True)
    add_textbox(s6b, 0.5, 0.62, 9.2, 0.3,
                "Proposition 4: every training exporter also exports inference — but not vice versa",
                font_size=12, color=C.textMuted, italic=True)

    # Nested circles
    outer = add_oval(s6b, 0.6, 1.15, 5.0, 3.3, hex_to_rgb("E0F7FA"))
    outer.line.color.rgb = C.accent
    outer.line.width = Pt(2.5)
    add_textbox(s6b, 0.8, 1.25, 2.5, 0.35, "Inference Exporters",
                font_size=12, color=C.accent, bold=True)

    inner = add_oval(s6b, 1.5, 1.75, 3.2, 2.2, hex_to_rgb("B2EBF2"))
    inner.line.color.rgb = hex_to_rgb("0288D1")
    inner.line.width = Pt(2.5)
    add_textbox(s6b, 2.1, 2.05, 2.0, 0.6, "Training\nExporters",
                font_size=13, color=hex_to_rgb("0288D1"), bold=True, align=PP_ALIGN.CENTER)

    # Labels
    for text, x, y in [("Canada", 2.3, 2.65), ("Kyrgyzstan†", 2.3, 2.93)]:
        add_textbox(s6b, x, y, 1.6, 0.28, text,
                    font_size=11, color=C.textDark, align=PP_ALIGN.CENTER)

    for text, x, y in [("Kosovo†", 0.8, 3.15), ("Finland", 0.8, 3.4),
                       ("Norway", 4.1, 3.15), ("UK", 4.1, 3.4)]:
        add_textbox(s6b, x, y, 1.2, 0.25, text,
                    font_size=10, color=C.accent, align=PP_ALIGN.CENTER)

    add_textbox(s6b, 2.7, 3.25, 0.6, 0.3, "EE",
                font_size=10, color=hex_to_rgb("0288D1"), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s6b, 1.6, 3.85, 0.5, 0.25, "IE",
                font_size=10, color=C.accent, bold=True, align=PP_ALIGN.CENTER)

    # Right cards
    add_card(s6b, 6.0, 1.15, 3.7, 1.65)
    add_rect(s6b, 6.0, 1.15, 0.07, 1.65, hex_to_rgb("0288D1"))
    add_textbox(s6b, 6.25, 1.2, 3.3, 0.3, "Why nesting holds",
                font_size=13, color=hex_to_rgb("0288D1"), bold=True)
    add_textbox(s6b, 6.25, 1.52, 3.3, 1.2,
                "Training has no distance penalty (τ = 0), so a training exporter wins the global cost competition.\n\nThat same cost advantage dominates the latency markup for nearby buyers → it also wins inference locally.",
                font_size=11, color=C.textDark)

    add_card(s6b, 6.0, 3.0, 3.7, 1.55)
    add_rect(s6b, 6.0, 3.0, 0.07, 1.55, C.warm)
    add_textbox(s6b, 6.25, 3.05, 3.3, 0.3, "But not vice versa",
                font_size=13, color=C.warm, bold=True)
    add_textbox(s6b, 6.25, 3.37, 3.3, 1.1,
                "An inference exporter only needs to beat neighbors within the latency threshold.\n\nIt may not be cheap enough to win the global training market, where every country competes head-to-head on cost alone.",
                font_size=11, color=C.textDark)

    # Bottom bar
    add_rect(s6b, 0.5, 4.75, 9.0, 0.5, C.darkBg)
    add_textbox(s6b, 0.7, 4.75, 8.6, 0.5,
                "Implication: regimes EI, ED, and DE cannot arise — the five equilibrium regimes in Proposition 1 are the only ones possible",
                font_size=12, color=C.textLight, bold=True, valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 15: COST RANKINGS (TABLE 3)
    # ============================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, C.lightBg)

    add_textbox(s7, 0.5, 0.25, 9.2, 0.5, "Raw Tariffs vs Cost-Recovery (Table 3)",
                font_name=HEADER, font_size=26, color=C.textDark, bold=True)
    add_textbox(s7, 0.5, 0.7, 9.2, 0.35,
                "Subsidized tariffs replaced with long-run marginal cost — 13 countries adjusted, 40 change trade regimes",
                font_size=12, color=C.textMuted, italic=True)

    rank_rows = [
        ["#", "(1) Raw Tariffs", "$/hr", "Type", "", "(2) Cost-Recovery", "$/hr", "Type"],
        ["1", "Iran", "$1.56", "DD*", "→", "Kyrgyzstan", "$1.58", "EE†"],
        ["2", "Turkmenistan", "$1.57", "DD*", "→", "Canada", "$1.59", "EE"],
        ["3", "Ethiopia", "$1.58", "EE†", "→", "Ethiopia", "$1.59", "II"],
        ["4", "Kyrgyzstan", "$1.58", "EE†", "→", "Kosovo", "$1.60", "IE†"],
        ["5", "Egypt", "$1.58", "EE†", "→", "Tajikistan", "$1.60", "II"],
        ["6", "Algeria", "$1.59", "EE†", "→", "Montenegro", "$1.60", "IE†"],
        ["7", "Canada", "$1.59", "EE", "→", "China", "$1.61", "II"],
        ["8", "Russia", "$1.59", "DD*", "→", "Ukraine", "$1.61", "IE†"],
        ["9", "South Africa", "$1.60", "II", "→", "Argentina", "$1.61", "II"],
        ["10", "Kosovo", "$1.60", "IE†", "→", "Colombia", "$1.61", "II"],
    ]

    tbl_shape = s7.shapes.add_table(len(rank_rows), 8, Inches(1.5), Inches(1.1), Inches(7.0), Inches(3.5))
    tbl = tbl_shape.table
    col_widths = [0.5, 1.5, 0.8, 0.7, 0.45, 1.5, 0.8, 0.7]
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = Inches(w)

    for i, row in enumerate(rank_rows):
        row_bg = C.darkBg if i == 0 else (hex_to_rgb("F8FAFC") if (i - 1) % 2 == 0 else C.cardBg)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j not in [1, 5] else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = BODY
                    if i == 0:
                        r.font.bold = True
                        r.font.color.rgb = C.textWhite
                    elif j in [3, 7]:
                        r.font.bold = True
                        r.font.color.rgb = regime_color(val)
                    elif j == 4:
                        r.font.color.rgb = C.textMuted
                    else:
                        r.font.color.rgb = C.textDark

    add_textbox(s7, 0.5, 4.9, 9.0, 0.3,
                "EE = training + inference exporter    IE = inference exporter    DD = domestic    II = full importer    * = sanctioned    † = developing",
                font_size=9, color=C.textMuted, align=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 16: SOVEREIGNTY RESHAPES EVERYTHING (TABLE 3b)
    # ============================================================
    s7b = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7b, C.lightBg)

    add_textbox(s7b, 0.5, 0.25, 9.2, 0.5, "Sovereignty Reshapes Everything (Table 3)",
                font_name=HEADER, font_size=26, color=C.textDark, bold=True)
    add_textbox(s7b, 0.5, 0.7, 9.2, 0.35,
                "Bilateral sovereignty premia shift most developing countries from exporters to domestic producers",
                font_size=12, color=C.textMuted, italic=True)

    sov_rows = [
        ["Country", "(1) Raw\nRank", "Type", "", "(2) Cost-Rec.\nRank", "Type", "", "(3) Bilateral\nRank", "Type"],
        ["Canada", "7", "EE", "▲", "2", "EE", "▲", "1", "EE"],
        ["Kyrgyzstan", "4", "EE†", "▲", "1", "EE†", "▼", "4", "EE†"],
        ["Iran", "1", "DD*", "▼", "24", "DD*", "—", "—", "DD*"],
        ["Ethiopia", "3", "EE†", "—", "3", "II", "▼", "8", "DD"],
        ["Kosovo", "10", "IE†", "▲", "4", "IE†", "▼", "9", "DD"],
        ["Finland", "24", "IE", "▲", "16", "IE", "▲", "6", "IE"],
        ["USA", "—", "II", "—", "—", "II", "▲", "2", "DD"],
        ["Turkmenistan", "2", "DD*", "▼", "19", "DD*", "—", "—", "DD*"],
        ["Norway", "21", "IE", "▲", "14", "IE", "▲", "5", "DD"],
    ]

    tbl_shape = s7b.shapes.add_table(len(sov_rows), 9, Inches(1.6), Inches(1.05), Inches(6.8), Inches(3.3))
    tbl = tbl_shape.table
    sov_col_widths = [1.4, 0.75, 0.65, 0.4, 0.85, 0.65, 0.4, 0.85, 0.65]
    for j, w in enumerate(sov_col_widths):
        tbl.columns[j].width = Inches(w)

    for i, row in enumerate(sov_rows):
        row_bg = C.darkBg if i == 0 else (hex_to_rgb("F8FAFC") if (i - 1) % 2 == 0 else C.cardBg)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(10 if i == 0 else 11)
                    r.font.name = BODY
                    if i == 0:
                        r.font.bold = True
                        r.font.color.rgb = C.textWhite
                    elif j in [2, 5, 8]:
                        r.font.bold = True
                        r.font.color.rgb = regime_color(val)
                    elif j in [3, 6]:
                        r.font.color.rgb = C.greenAcc if val == "▲" else (C.redAcc if val == "▼" else C.textMuted)
                    else:
                        r.font.color.rgb = C.textDark

    add_textbox(s7b, 0.5, 4.55, 9.0, 0.25,
                "— = not in top 25 (USA high-cost importer) or no bilateral rank (Iran, Turkmenistan sanctioned → forced domestic)",
                font_size=9, color=C.textMuted, italic=True, align=PP_ALIGN.CENTER)

    add_rect(s7b, 0.5, 4.75, 4.2, 0.5, C.darkBg)
    add_textbox(s7b, 0.7, 4.75, 3.8, 0.5,
                "Only Canada + Kyrgyzstan remain exporters under bilateral spec",
                font_size=11, color=C.textLight, valign=MSO_ANCHOR.MIDDLE)
    add_rect(s7b, 4.9, 4.75, 4.6, 0.5, C.darkBg)
    add_textbox(s7b, 5.1, 4.75, 4.2, 0.5,
                "Trusted advanced economies rise as developing exporters fall",
                font_size=11, color=C.textLight, valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 17: SOVEREIGNTY PREMIUM — CENTRAL FINDING
    # ============================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, C.darkBg)

    add_rect(s8, 0, 0, 10, 0.06, C.warm)
    add_textbox(s8, 1.4, 0.5, 7, 0.55, "The Central Finding",
                font_name=HEADER, font_size=30, color=C.textWhite, bold=True)

    add_textbox(s8, 0.8, 1.5, 8.4, 1.2,
                "Countries best positioned to export compute on cost grounds are precisely those excluded by bilateral trust deficits.",
                font_name=HEADER, font_size=24, color=C.warm, italic=True)

    consequences = [
        {"title": "Cost advantage exists",
         "desc": "Developing countries with cheap energy rank among the lowest-cost producers under cost-recovery pricing.",
         "accent": C.greenAcc},
        {"title": "Trust deficits eliminate it",
         "desc": "Bilateral sovereignty premia shift nearly all developing countries to domestic production or importing.",
         "accent": C.redAcc},
        {"title": "Welfare cost: ~1.5%",
         "desc": "Demand-weighted welfare loss from sovereignty premia is 1.5% of compute spending — modest but nontrivial.",
         "accent": C.warm},
    ]

    for i, c in enumerate(consequences):
        xPos = 0.7 + i * 3.1
        add_rect(s8, xPos, 3.0, 2.85, 2.1, C.medBg)
        add_rect(s8, xPos, 3.0, 2.85, 0.05, c["accent"])
        add_textbox(s8, xPos + 0.2, 3.2, 2.45, 0.4, c["title"],
                    font_size=14, color=c["accent"], bold=True)
        add_textbox(s8, xPos + 0.2, 3.65, 2.45, 1.3, c["desc"],
                    font_size=12, color=C.textLight)

    # ============================================================
    # SLIDE 18: DUAL IMPLICATION
    # ============================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, C.lightBg)

    add_textbox(s9, 0.7, 0.4, 9, 0.7, "A Dual Implication",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)
    add_textbox(s9, 0.7, 1.0, 9, 0.35, "The narrow cost spread cuts both ways",
                font_size=15, color=C.textMuted, italic=True)

    # Left card
    add_card(s9, 0.7, 1.7, 4.15, 3.3)
    add_rect(s9, 0.7, 1.7, 0.07, 3.3, C.greenAcc)
    add_textbox(s9, 1.1, 1.9, 3.5, 0.4, "Easiest Sector to Enter",
                font_name=HEADER, font_size=18, color=C.greenAcc, bold=True)
    add_textbox(s9, 1.1, 2.5, 3.5, 2.2,
                "Because hardware is globally priced and dominates costs, the absolute barrier to cost-competitiveness is low. Countries with cheap electricity and favorable cooling face a small gap to close.",
                font_size=14, color=C.textDark)

    # Right card
    add_card(s9, 5.15, 1.7, 4.15, 3.3)
    add_rect(s9, 5.15, 1.7, 0.07, 3.3, C.redAcc)
    add_textbox(s9, 5.5, 1.9, 3.5, 0.4, "Most Vulnerable to Frictions",
                font_name=HEADER, font_size=18, color=C.redAcc, bold=True)
    add_textbox(s9, 5.5, 2.5, 3.5, 2.2,
                "The same narrowness means even modest frictions — a sovereignty premium, a governance penalty, or a higher cost of capital — are sufficient to alter a country's trade regime entirely.",
                font_size=14, color=C.textDark)

    # ============================================================
    # SLIDE 19: SENSITIVITY ANALYSIS
    # ============================================================
    s9b = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9b, C.lightBg)

    add_textbox(s9b, 0.5, 0.2, 9.2, 0.45, "Robustness: Sensitivity Analysis",
                font_name=HEADER, font_size=25, color=C.textDark, bold=True)
    add_textbox(s9b, 0.5, 0.6, 9.2, 0.3, "Cost rankings are stable across plausible parameter variation",
                font_size=12, color=C.textMuted, italic=True)

    robCards = [
        {"title": "Hardware cost share (ρ ± 4%)",
         "finding": "Higher ρ compresses the locally-penalized share → governance penalties muted → developing countries gain competitiveness.",
         "detail": "Lower ρ widens the electricity and construction share → cold-climate, low-cost countries strengthen their advantage.",
         "color": C.accent},
        {"title": "Uniform sovereignty premium (λ = 10%)",
         "finding": "Most countries produce domestically under either bilateral or uniform specification — regime assignments are similar.",
         "detail": "Only 2 countries (Canada, Kyrgyzstan) remain exporters under both. The narrow cost spread means 10% is enough to flip regimes.",
         "color": C.warm},
        {"title": "Sovereignty tiers (3 workload segments)",
         "finding": "Sovereign workloads (10% of demand): domestic only. Regulated (20%): higher compatibility weight. Commercial (70%): standard.",
         "detail": "Tiered approach more realistic than uniform — governments rarely apply the same standard to classified vs. commercial AI.",
         "color": C.accent},
    ]

    for i, c in enumerate(robCards):
        yPos = 1.05 + i * 1.25
        add_card(s9b, 0.5, yPos, 9.0, 1.1)
        add_rect(s9b, 0.5, yPos, 0.07, 1.1, c["color"])
        add_textbox(s9b, 0.8, yPos + 0.08, 8.5, 0.3, c["title"],
                    font_size=14, color=c["color"], bold=True)
        add_textbox(s9b, 0.8, yPos + 0.38, 8.5, 0.3, c["finding"],
                    font_size=12, color=C.textDark)
        add_textbox(s9b, 0.8, yPos + 0.7, 8.5, 0.3, c["detail"],
                    font_size=10, color=C.textMuted)

    add_rect(s9b, 0.5, 4.75, 9.0, 0.5, C.darkBg)
    add_textbox(s9b, 0.7, 4.75, 8.6, 0.5,
                "Core result stable: the binding constraint is institutional credibility, not energy cost, across all specifications",
                font_size=12, color=C.textLight, bold=True, valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 20: CAVEATS & EXTENSIONS
    # ============================================================
    s9c = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9c, C.lightBg)

    add_textbox(s9c, 0.5, 0.2, 9.2, 0.45, "Caveats & Extensions",
                font_name=HEADER, font_size=25, color=C.textDark, bold=True)
    add_textbox(s9c, 0.5, 0.6, 9.2, 0.3,
                "Factors outside the current model that would further narrow the set of viable exporters",
                font_size=12, color=C.textMuted, italic=True)

    extCards = [
        {"title": "GPU Export Controls",
         "body": "US restrictions bar Iran, Russia, Belarus, and (partially) China from acquiring frontier GPUs. Model's cost advantage for sanctioned countries is moot without hardware access.",
         "x": 0.5, "y": 1.05},
        {"title": "Endogenous Electricity Prices",
         "body": "A 100 MW data center in Kyrgyzstan would consume ~10% of national generation. At scale, marginal cost pricing would push electricity above the current average tariff.",
         "x": 5.1, "y": 1.05},
        {"title": "Cost of Capital",
         "body": "Calibration assumes uniform WACC. In practice, an OECD hyperscaler at 8% faces $1.58/hr; a developing-country entrant at 15% faces $1.62/hr — erasing cost advantage from cheap energy.",
         "x": 0.5, "y": 2.75},
        {"title": "Next-Gen GPUs (B200+)",
         "body": "Successor GPUs deliver ~4× training throughput at ~1 kW. Higher power draw widens the absolute electricity gap, modestly strengthening developing-country advantage on energy costs.",
         "x": 5.1, "y": 2.75},
    ]

    for ec in extCards:
        add_card(s9c, ec["x"], ec["y"], 4.4, 1.5)
        add_rect(s9c, ec["x"], ec["y"], 4.4, 0.06, C.accent)
        add_textbox(s9c, ec["x"] + 0.2, ec["y"] + 0.15, 4.0, 0.35, ec["title"],
                    font_size=13, color=C.textDark, bold=True)
        add_textbox(s9c, ec["x"] + 0.2, ec["y"] + 0.55, 4.0, 0.85, ec["body"],
                    font_size=11, color=C.textMuted)

    add_rect(s9c, 0.5, 4.55, 9.0, 0.7, C.darkBg)
    add_rich_text(s9c, 0.7, 4.55, 8.6, 0.7, [
        {"text": "Direction of bias: ", "color": C.warm, "bold": True, "font_size": 12},
        {"text": "all four factors narrow the set of viable exporters beyond the baseline calibration. GPU controls and capital costs work against developing countries; endogenous prices work against small countries; GPU upgrades slightly favor energy-rich ones.",
         "color": C.textLight, "font_size": 11},
    ], valign=MSO_ANCHOR.MIDDLE)

    # ============================================================
    # SLIDE 21: POLICY IMPLICATIONS
    # ============================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, C.lightBg)

    add_textbox(s10, 0.7, 0.4, 9, 0.7, "Policy Implications",
                font_name=HEADER, font_size=32, color=C.textDark, bold=True)

    policies = [
        {"title": "Training vs Inference", "desc": "Restricting training imports raises costs without offsetting proximity gains. Inference has a genuine latency-based case for domestic production.", "x": 0.7, "y": 1.3},
        {"title": "Institutional Credibility", "desc": "The binding constraint is not energy cost but institutional trust: non-sanctioned status, credible contracts, network connectivity, regulatory stability.", "x": 5.15, "y": 1.3},
        {"title": "Data Localization Risk", "desc": "Broad data localization requirements risk foreclosing both import savings and regional export opportunities — the very gains the model predicts.", "x": 0.7, "y": 3.1},
        {"title": "Resource Curse Parallels", "desc": "Concentrated compute export revenues could trigger Dutch disease dynamics. Revenue-sharing models (sovereign wealth fund vs elite capture) determine outcomes.", "x": 5.15, "y": 3.1},
    ]

    for p in policies:
        add_card(s10, p["x"], p["y"], 4.15, 1.55)
        add_textbox(s10, p["x"] + 0.25, p["y"] + 0.15, 3.6, 0.4, p["title"],
                    font_name=HEADER, font_size=15, color=C.textDark, bold=True)
        add_textbox(s10, p["x"] + 0.25, p["y"] + 0.6, 3.6, 0.85, p["desc"],
                    font_size=12, color=C.textMuted)

    # ============================================================
    # SLIDE 22: KEY TAKEAWAY / CLOSING
    # ============================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11, C.darkBg)

    add_rect(s11, 0, 0, 10, 0.06, C.accent)
    add_textbox(s11, 0.8, 0.6, 8, 0.6, "Key Takeaway",
                font_name=HEADER, font_size=30, color=C.textWhite, bold=True)

    add_textbox(s11, 0.8, 1.5, 8.4, 1.2,
                "Compute is the easiest sector for developing countries to enter on cost grounds — and the most vulnerable to the small frictions that keep them out.",
                font_name=HEADER, font_size=22, color=C.accent)

    takeaways = [
        "Hardware dominance compresses cross-country cost spread to 12–20%",
        "Cost-recovery pricing reveals genuine comparative advantage in energy-rich countries",
        "Bilateral sovereignty premia eliminate developing-country exports",
        "Institutional credibility, not energy cost, is the binding constraint",
    ]
    for i, t in enumerate(takeaways):
        yPos = 3.0 + i * 0.55
        add_textbox(s11, 0.8, yPos + 0.05, 0.25, 0.25, "→",
                    font_size=14, color=C.accent, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(s11, 1.2, yPos, 8, 0.45, t,
                    font_size=15, color=C.textLight)

    # Footer
    add_rect(s11, 0, 5.2, 10, 0.425, C.medBg)
    add_textbox(s11, 0.8, 5.25, 8, 0.35,
                "Michael Lokshin  •  March 2026  •  michael.lokshin@gmail.com",
                font_size=11, color=C.textMuted)

    # ============================================================
    # SAVE
    # ============================================================
    output_dir = os.path.dirname(os.path.abspath(__file__))
    doc_dir = os.path.join(os.path.dirname(output_dir), "Documents")
    output_path = os.path.join(doc_dir, "flop_trade_presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
